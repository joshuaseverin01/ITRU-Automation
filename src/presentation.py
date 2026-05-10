"""Presentation draft generation, validation, and PowerPoint export helpers."""

from __future__ import annotations

from dataclasses import asdict, dataclass, field
from io import BytesIO
import json
import re
from typing import Any, Mapping, Sequence
from urllib import error, request

import pandas as pd


PRESENTATION_PURPOSES = {
    "executive_summary": "Executive summary",
    "product_demo": "Product demo",
    "investor_update": "Investor update",
    "sales_deck": "Sales deck",
    "strategy_memo": "Strategy memo",
    "research_presentation": "Research presentation",
    "public_explainer": "Public-facing explainer",
}
PRESENTATION_AUDIENCES = {
    "executives": "Executives",
    "customers": "Customers",
    "investors": "Investors",
    "internal_team": "Internal team",
    "technical_team": "Technical team",
    "public_readers": "Public readers",
}
PRESENTATION_STYLES = {
    "concise": "Concise",
    "data_heavy": "Data-heavy",
    "visual_first": "Visual-first",
    "narrative": "Narrative",
    "consulting_style": "Consulting-style",
}
SOURCE_TYPES = {"topic", "blog_post", "source_material"}
SLIDE_COUNTS = {5, 8, 10, 12}
SLIDE_TYPES = {
    "title",
    "agenda",
    "context",
    "problem",
    "insight",
    "data",
    "comparison",
    "process",
    "recommendation",
    "closing",
}
VISUAL_TYPES = {
    "none",
    "chart_placeholder",
    "table_placeholder",
    "image_placeholder",
    "metric_cards",
    "two_column_comparison",
    "process_flow",
    "map_placeholder",
}


@dataclass
class PresentationSlide:
    """One validated slide in a generated presentation deck."""

    id: str
    slideNumber: int
    type: str
    title: str
    subtitle: str | None = None
    bullets: list[str] = field(default_factory=list)
    takeaway: str | None = None
    speakerNotes: str | None = None
    visualType: str = "none"
    visualPrompt: str | None = None
    chartData: list[dict[str, str | int | float]] = field(default_factory=list)
    tableData: list[list[str]] = field(default_factory=list)


@dataclass
class PresentationDeck:
    """Validated presentation deck data ready for preview and PPTX export."""

    deckTitle: str
    audience: str
    purpose: str
    style: str
    slides: list[PresentationSlide]
    deckSubtitle: str | None = None
    sourceSummary: str | None = None


@dataclass(frozen=True)
class PresentationGenerationResult:
    """Deck generation result with provenance for UI messaging."""

    deck: PresentationDeck
    generation_mode: str
    message: str | None = None


def normalize_presentation_input(raw_input: Mapping[str, Any]) -> dict[str, Any]:
    """Validate and normalize presentation-generation input from the UI."""

    source_type = str(raw_input.get("sourceType") or "topic").strip()
    if source_type not in SOURCE_TYPES:
        raise ValueError("Choose a valid presentation source type.")

    slide_count = int(raw_input.get("slideCount") or 8)
    if slide_count not in SLIDE_COUNTS:
        raise ValueError("Slide count must be 5, 8, 10, or 12.")

    normalized = {
        "sourceType": source_type,
        "topic": _clean_text(raw_input.get("topic")),
        "sourceText": _clean_text(raw_input.get("sourceText")),
        "tone": _clean_text(raw_input.get("tone")) or "Executive strategy language",
        "purpose": _choice(raw_input.get("purpose"), PRESENTATION_PURPOSES, "executive_summary"),
        "audience": _choice(raw_input.get("audience"), PRESENTATION_AUDIENCES, "executives"),
        "style": _choice(raw_input.get("style"), PRESENTATION_STYLES, "consulting_style"),
        "slideCount": slide_count,
        "includeSpeakerNotes": bool(raw_input.get("includeSpeakerNotes", True)),
        "includeVisualSuggestions": bool(raw_input.get("includeVisualSuggestions", True)),
    }
    source_material = _source_material(normalized)
    if not source_material.strip():
        raise ValueError("Add a topic, source material, or generated blog draft before creating a presentation.")
    return normalized


def build_presentation_deck_draft(
    raw_input: Mapping[str, Any],
    *,
    zone_df: pd.DataFrame | None = None,
    monthly_df: pd.DataFrame | None = None,
    ai_api_key: str | None = None,
    ai_model: str | None = None,
    prefer_ai: bool = True,
) -> PresentationGenerationResult:
    """Generate a structured presentation deck, using AI when configured and falling back locally."""

    normalized = normalize_presentation_input(raw_input)
    if prefer_ai and ai_api_key:
        try:
            deck_data = _generate_deck_with_openai(normalized, ai_api_key=ai_api_key, ai_model=ai_model)
            deck = validate_presentation_deck(deck_data)
            return PresentationGenerationResult(deck=deck, generation_mode="AI", message="Generated with the configured AI model.")
        except Exception as exc:
            fallback = _build_deterministic_deck(normalized, zone_df=zone_df, monthly_df=monthly_df)
            return PresentationGenerationResult(
                deck=fallback,
                generation_mode="Deterministic fallback",
                message=f"AI generation was unavailable, so a local structured draft was created instead: {exc}",
            )

    return PresentationGenerationResult(
        deck=_build_deterministic_deck(normalized, zone_df=zone_df, monthly_df=monthly_df),
        generation_mode="Deterministic local draft",
        message="No AI API key was configured, so a deterministic structured draft was created from the active inputs.",
    )


def validate_presentation_deck(raw_deck: Mapping[str, Any] | PresentationDeck) -> PresentationDeck:
    """Validate a presentation deck and repair safe metadata such as ids and slide numbers."""

    if isinstance(raw_deck, PresentationDeck):
        raw_deck = presentation_deck_to_dict(raw_deck)

    deck_title = _truncate(_required_text(raw_deck.get("deckTitle"), "deckTitle"), 140)
    slides_raw = raw_deck.get("slides")
    if not isinstance(slides_raw, Sequence) or isinstance(slides_raw, (str, bytes)):
        raise ValueError("Presentation deck must contain a slides array.")
    if not 3 <= len(slides_raw) <= 15:
        raise ValueError("Presentation deck must contain between 3 and 15 slides.")

    slides: list[PresentationSlide] = []
    seen_ids: set[str] = set()
    for index, raw_slide in enumerate(slides_raw, start=1):
        if not isinstance(raw_slide, Mapping):
            raise ValueError(f"Slide {index} must be an object.")
        slide_id = _clean_text(raw_slide.get("id")) or f"slide-{index}"
        if slide_id in seen_ids:
            slide_id = f"{slide_id}-{index}"
        seen_ids.add(slide_id)
        slide_type = _choice(raw_slide.get("type"), {key: key for key in SLIDE_TYPES}, "insight")
        visual_type = _choice(raw_slide.get("visualType"), {key: key for key in VISUAL_TYPES}, "none")
        slide = PresentationSlide(
            id=slide_id,
            slideNumber=index,
            type=slide_type,
            title=_truncate(_required_text(raw_slide.get("title"), f"slide {index} title"), 140),
            subtitle=_optional_text(raw_slide.get("subtitle"), 200),
            bullets=_list_of_text(raw_slide.get("bullets"), max_items=6, max_chars=180),
            takeaway=_optional_text(raw_slide.get("takeaway"), 220),
            speakerNotes=_optional_text(raw_slide.get("speakerNotes"), 1200),
            visualType=visual_type,
            visualPrompt=_optional_text(raw_slide.get("visualPrompt"), 400),
            chartData=_chart_rows(raw_slide.get("chartData")),
            tableData=_table_rows(raw_slide.get("tableData")),
        )
        slides.append(slide)

    return PresentationDeck(
        deckTitle=deck_title,
        deckSubtitle=_optional_text(raw_deck.get("deckSubtitle"), 220),
        audience=_choice(raw_deck.get("audience"), PRESENTATION_AUDIENCES, "executives"),
        purpose=_choice(raw_deck.get("purpose"), PRESENTATION_PURPOSES, "executive_summary"),
        style=_choice(raw_deck.get("style"), PRESENTATION_STYLES, "consulting_style"),
        sourceSummary=_optional_text(raw_deck.get("sourceSummary"), 700),
        slides=slides,
    )


def presentation_deck_to_dict(deck: PresentationDeck) -> dict[str, Any]:
    """Return a JSON-serializable representation of a presentation deck."""

    return asdict(deck)


def presentation_deck_to_pptx_bytes(deck: PresentationDeck | Mapping[str, Any]) -> bytes:
    """Export a validated deck to a PowerPoint file."""

    deck_obj = validate_presentation_deck(deck) if not isinstance(deck, PresentationDeck) else deck
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:  # pragma: no cover - exercised in deployment if dependency missing
        raise RuntimeError("PowerPoint export requires python-pptx. Install requirements.txt and try again.") from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    colors = {
        "dark": RGBColor(17, 24, 39),
        "muted": RGBColor(107, 114, 128),
        "border": RGBColor(183, 228, 183),
        "accent": RGBColor(28, 181, 28),
        "light": RGBColor(248, 255, 248),
        "white": RGBColor(255, 255, 255),
    }

    for slide_data in deck_obj.slides:
        slide = prs.slides.add_slide(blank_layout)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = colors["light"]

        if slide_data.type == "title":
            _add_textbox(slide, slide_data.title, 0.85, 1.55, 11.6, 1.35, Pt(36), colors["dark"], bold=True)
            if slide_data.subtitle:
                _add_textbox(slide, slide_data.subtitle, 0.9, 2.85, 10.8, 0.55, Pt(17), colors["muted"])
            if deck_obj.sourceSummary:
                _add_textbox(slide, deck_obj.sourceSummary, 0.9, 4.35, 10.8, 0.75, Pt(13), colors["muted"])
            _add_accent_bar(slide, 0.9, 5.55, 2.7, 0.08, colors["accent"])
        else:
            _add_textbox(slide, slide_data.title, 0.55, 0.42, 8.35, 0.65, Pt(23), colors["dark"], bold=True)
            if slide_data.subtitle:
                _add_textbox(slide, slide_data.subtitle, 0.58, 1.08, 8.45, 0.42, Pt(11), colors["muted"])
            _add_slide_body(slide, slide_data, colors, MSO_SHAPE, Inches, Pt)

        _add_footer(slide, deck_obj.deckTitle, slide_data.slideNumber, len(deck_obj.slides), colors, Inches, Pt)
        if slide_data.speakerNotes:
            _add_speaker_notes(slide, slide_data.speakerNotes)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def build_presentation_system_prompt() -> str:
    """Prompt instructions for AI structured presentation generation."""

    return (
        "You are an expert presentation strategist and product/content editor. "
        "Transform user-provided material into a structured slide deck that can be rendered programmatically. "
        "Return only valid JSON matching the requested schema. Do not return markdown or commentary. "
        "Create slide-ready content, not long prose. Use clear insight-driven slide titles. "
        "Keep bullets concise. Do not fabricate facts, metrics, dates, company names, quotes, or results."
    )


def build_presentation_user_prompt(normalized_input: Mapping[str, Any]) -> str:
    """Build the user prompt sent to the AI provider."""

    return "\n".join(
        [
            "Create a presentation deck using the following settings:",
            "",
            f"Source type: {normalized_input['sourceType']}",
            f"Purpose: {normalized_input['purpose']}",
            f"Audience: {normalized_input['audience']}",
            f"Style: {normalized_input['style']}",
            f"Tone: {normalized_input['tone']}",
            f"Target slide count: {normalized_input['slideCount']}",
            f"Include speaker notes: {normalized_input['includeSpeakerNotes']}",
            f"Include visual suggestions: {normalized_input['includeVisualSuggestions']}",
            "",
            "Source material:",
            _source_material(normalized_input),
            "",
            "Return JSON with keys: deckTitle, deckSubtitle, audience, purpose, style, sourceSummary, slides.",
            "Each slide must include id, slideNumber, type, title, bullets, takeaway, speakerNotes, visualType, and visualPrompt where relevant.",
            f"Return exactly {normalized_input['slideCount']} slides unless the content is too thin. If too thin, return at least 3 slides.",
            "Use slideNumber values starting at 1. The first slide should be a title slide. The final slide should be a closing, recommendation, or takeaway slide.",
            "Do not include unsupported keys. Do not copy long paragraphs into bullets. Do not fabricate metrics.",
        ]
    )


def _generate_deck_with_openai(normalized_input: Mapping[str, Any], *, ai_api_key: str, ai_model: str | None) -> dict[str, Any]:
    # Keep this dependency-free for Streamlit Cloud portability. Validation remains the source of truth.
    payload = {
        "model": ai_model or "gpt-4o-mini",
        "messages": [
            {"role": "system", "content": build_presentation_system_prompt()},
            {"role": "user", "content": build_presentation_user_prompt(normalized_input)},
        ],
        "temperature": 0.35,
        "response_format": {"type": "json_object"},
    }
    api_request = request.Request(
        "https://api.openai.com/v1/chat/completions",
        data=json.dumps(payload).encode("utf-8"),
        headers={"Authorization": f"Bearer {ai_api_key}", "Content-Type": "application/json"},
        method="POST",
    )
    try:
        with request.urlopen(api_request, timeout=60) as response:
            response_data = json.loads(response.read().decode("utf-8"))
    except error.HTTPError as exc:
        detail = exc.read().decode("utf-8", errors="ignore")
        raise RuntimeError(f"AI request failed with status {exc.code}: {detail[:240]}") from exc

    content = response_data["choices"][0]["message"]["content"]
    return json.loads(content)


def _build_deterministic_deck(
    normalized_input: Mapping[str, Any],
    *,
    zone_df: pd.DataFrame | None,
    monthly_df: pd.DataFrame | None,
) -> PresentationDeck:
    source = _source_material(normalized_input)
    topic = normalized_input.get("topic") or _infer_title_from_source(source) or "Flexworks market intelligence"
    purpose = str(normalized_input["purpose"])
    audience = str(normalized_input["audience"])
    style = str(normalized_input["style"])
    slide_count = int(normalized_input["slideCount"])
    include_notes = bool(normalized_input["includeSpeakerNotes"])
    include_visuals = bool(normalized_input["includeVisualSuggestions"])

    analysis = _summarize_zone_data(zone_df)
    timing = _summarize_monthly_data(monthly_df)
    title = _presentation_title(topic, analysis)
    source_summary = _truncate(_plain_text(source), 680)
    slides = _deterministic_slide_plan(
        title=title,
        topic=str(topic),
        analysis=analysis,
        timing=timing,
        slide_count=slide_count,
        include_notes=include_notes,
        include_visuals=include_visuals,
    )
    return validate_presentation_deck(
        {
            "deckTitle": title,
            "deckSubtitle": "Structured presentation draft generated from Flexworks analysis inputs",
            "audience": audience,
            "purpose": purpose,
            "style": style,
            "sourceSummary": source_summary,
            "slides": slides,
        }
    )


def _deterministic_slide_plan(
    *,
    title: str,
    topic: str,
    analysis: Mapping[str, Any],
    timing: Mapping[str, Any],
    slide_count: int,
    include_notes: bool,
    include_visuals: bool,
) -> list[dict[str, Any]]:
    top_zones = analysis.get("top_zones") or []
    bottom_zones = analysis.get("bottom_zones") or []
    metric_label = analysis.get("metric_label") or "selected metric"
    spread_text = analysis.get("spread_text") or "not available"
    top_phrase = ", ".join(top_zones[:3]) if top_zones else "the leading modeled zones"
    bottom_phrase = ", ".join(bottom_zones[:3]) if bottom_zones else "the lower-ranked modeled zones"
    timing_line = timing.get("summary") or "Review monthly or timestamp-level outputs before making timing claims."

    candidates = [
        {
            "type": "title",
            "title": title,
            "subtitle": f"Draft narrative for {topic}",
            "bullets": [],
            "takeaway": "Use this deck as a reviewable first draft, not a final published claim.",
            "visualType": "none",
        },
        {
            "type": "context",
            "title": "Battery revenue strategy starts with location",
            "bullets": [
                "Local market conditions can change modeled arbitrage value by zone.",
                "Zone-level analysis helps teams prioritize where deeper diligence should start.",
                "Flexworks turns dense simulation exports into a repeatable market-intelligence workflow.",
            ],
            "takeaway": "The deck frames location as a first-order strategy variable.",
            "visualType": "map_placeholder",
            "visualPrompt": "ISO zone map with revenue intensity by zone.",
        },
        {
            "type": "data",
            "title": f"Top zones separate on {metric_label}",
            "bullets": [
                f"Highest-ranked zones: {top_phrase}.",
                f"Lowest-ranked zones: {bottom_phrase}.",
                f"Top-to-bottom spread: {spread_text}.",
            ],
            "takeaway": "The modeled spread indicates where location may materially affect value.",
            "visualType": "table_placeholder",
            "visualPrompt": "Ranked table with top and bottom zones and selected metric values.",
        },
        {
            "type": "insight",
            "title": "The result is a screening signal, not a final siting answer",
            "bullets": [
                "Modeled leaders should move into deeper diligence, not automatic selection.",
                "Low-ranked zones may still work with different operating objectives or value stacks.",
                "The next question is whether the spread is persistent across time.",
            ],
            "takeaway": "Use the output to focus follow-up analysis where it matters most.",
            "visualType": "two_column_comparison",
            "visualPrompt": "Two-column comparison of screening signal versus final diligence needs.",
        },
        {
            "type": "data",
            "title": "Timing determines whether revenue is persistent or event-driven",
            "bullets": [
                timing_line,
                "Monthly patterns help distinguish steady zone advantage from a small number of peak periods.",
                "Persistent leaders may support stronger planning confidence than one-off spikes.",
            ],
            "takeaway": "Time-series review helps separate durable signals from event-driven results.",
            "visualType": "chart_placeholder",
            "visualPrompt": "Monthly revenue trend by zone or category.",
        },
        {
            "type": "process",
            "title": "Flexworks compresses the manual workflow into a repeatable loop",
            "bullets": [
                "Clean exports and normalize device or node records.",
                "Join devices to zones and aggregate performance metrics.",
                "Generate maps, rankings, animations, exports, and narrative drafts.",
            ],
            "takeaway": "The value is speed plus traceability from raw export to strategy output.",
            "visualType": "process_flow",
            "visualPrompt": "Process flow from simulation export to analysis, visualization, and presentation.",
        },
        {
            "type": "recommendation",
            "title": "Use zone rankings to prioritize the next diligence queue",
            "bullets": [
                "Move top zones into deeper operational and market-rule review.",
                "Pressure-test weaker zones for value stacking or alternative use cases.",
                "Compare results against asset specs, dispatch assumptions, and customer priorities.",
            ],
            "takeaway": "The recommended action is focused follow-up, not one-size-fits-all deployment.",
            "visualType": "metric_cards",
            "visualPrompt": "Three action cards: prioritize, pressure-test, validate.",
        },
        {
            "type": "closing",
            "title": "The strategic question is where flexibility earns the most",
            "bullets": [
                "Zone-level modeling makes locational variation visible.",
                "Temporal views show how performance evolves across the modeled period.",
                "Exports turn analysis into a reusable customer or investment conversation.",
            ],
            "takeaway": "Flexworks helps translate simulation outputs into market strategy.",
            "visualType": "none",
        },
        {
            "type": "comparison",
            "title": "High-value zones deserve a different conversation than low-value zones",
            "bullets": [
                "Top zones can anchor opportunity sizing and partnership discussions.",
                "Middle zones may depend more on operations, timing, or portfolio effects.",
                "Bottom zones require stronger justification before advancing.",
            ],
            "takeaway": "The same battery strategy may not travel cleanly across every zone.",
            "visualType": "two_column_comparison",
        },
        {
            "type": "data",
            "title": "The export package supports both analysis and communication",
            "bullets": [
                "Cleaned CSVs support auditability and additional modeling.",
                "Maps and animations make locational differences easier to explain.",
                "Draft narratives help teams move faster from analysis to decision support.",
            ],
            "takeaway": "A repeatable output package makes the work easier to review and share.",
            "visualType": "chart_placeholder",
        },
        {
            "type": "recommendation",
            "title": "Before publication, confirm assumptions and market context",
            "bullets": [
                "Validate final asset specifications and dispatch assumptions.",
                "Check whether market-rule or event context changes the interpretation.",
                "Review claims before using the deck externally.",
            ],
            "takeaway": "The draft should be reviewed before external publication.",
            "visualType": "none",
        },
        {
            "type": "closing",
            "title": "From simulation output to investment-ready market intelligence",
            "bullets": [
                "The workflow reduces manual analysis effort.",
                "The deck captures the main implications of the modeled results.",
                "The next step is targeted diligence on the strongest and weakest zones.",
            ],
            "takeaway": "Location, timing, and value stacking should guide the next analysis loop.",
            "visualType": "none",
        },
    ]

    selected = candidates[:slide_count]
    if selected[-1]["type"] != "closing":
        selected[-1] = candidates[7]

    for index, slide in enumerate(selected, start=1):
        slide["id"] = f"slide-{index}"
        slide["slideNumber"] = index
        if not include_notes:
            slide["speakerNotes"] = None
        else:
            slide["speakerNotes"] = _speaker_note_for(slide)
        if not include_visuals:
            slide["visualType"] = "none"
            slide["visualPrompt"] = None
        else:
            slide.setdefault("visualPrompt", None)
    return selected


def _summarize_zone_data(zone_df: pd.DataFrame | None) -> dict[str, Any]:
    if zone_df is None or zone_df.empty:
        return {"metric_label": "selected metric", "top_zones": [], "bottom_zones": [], "spread_text": "not available"}
    metric_column = next(
        (
            column
            for column in ["Selected_Metric", "Revenue_per_kW", "Annualized_Revenue", "Opportunity_Score", "Risk_Adjusted_Score"]
            if column in zone_df.columns and not pd.to_numeric(zone_df[column], errors="coerce").dropna().empty
        ),
        None,
    )
    zone_column = next((column for column in ["Zone", "Node_ID", "Device_ID", "Device", "Node_Name"] if column in zone_df.columns), None)
    if metric_column is None or zone_column is None:
        return {"metric_label": "selected metric", "top_zones": [], "bottom_zones": [], "spread_text": "not available"}
    working = zone_df[[zone_column, metric_column]].copy()
    working[metric_column] = pd.to_numeric(working[metric_column], errors="coerce")
    working[zone_column] = working[zone_column].astype(str).str.strip()
    working = working.dropna(subset=[zone_column, metric_column])
    working = working.loc[working[zone_column] != ""].copy()
    if working.empty:
        return {"metric_label": _label(metric_column), "top_zones": [], "bottom_zones": [], "spread_text": "not available"}
    grouped = working.groupby(zone_column, as_index=False)[metric_column].mean().sort_values(metric_column, ascending=False)
    spread = float(grouped[metric_column].max() - grouped[metric_column].min())
    return {
        "metric_label": _label(metric_column),
        "top_zones": grouped.head(3)[zone_column].astype(str).tolist(),
        "bottom_zones": grouped.tail(3).sort_values(metric_column, ascending=True)[zone_column].astype(str).tolist(),
        "spread_text": _format_value(spread, metric_column),
    }


def _summarize_monthly_data(monthly_df: pd.DataFrame | None) -> dict[str, Any]:
    if monthly_df is None or monthly_df.empty or "Revenue" not in monthly_df.columns:
        return {"summary": "Review note: add time-series interpretation after monthly patterns are reviewed."}
    time_column = next((column for column in ["Month", "Timestamp", "Time"] if column in monthly_df.columns), None)
    if time_column is None:
        return {"summary": "Review note: add time-series interpretation after monthly patterns are reviewed."}
    working = monthly_df.copy()
    working["Revenue"] = pd.to_numeric(working["Revenue"], errors="coerce")
    working["Time"] = pd.to_datetime(working[time_column], errors="coerce")
    working = working.dropna(subset=["Revenue", "Time"])
    if working.empty:
        return {"summary": "Review note: add time-series interpretation after monthly patterns are reviewed."}
    by_time = working.groupby("Time", as_index=False)["Revenue"].sum().sort_values("Time")
    peak = by_time.sort_values("Revenue", ascending=False).iloc[0]
    period_count = len(by_time)
    return {
        "summary": (
            f"The time-series data covers {period_count} period(s), with the highest modeled revenue in "
            f"{pd.to_datetime(peak['Time']).strftime('%B %Y')} at {_format_value(float(peak['Revenue']), 'Annualized_Revenue')}."
        )
    }


def _presentation_title(topic: object, analysis: Mapping[str, Any]) -> str:
    top_zones = analysis.get("top_zones") or []
    if top_zones:
        return "Turning Zone-Level Battery Revenue Into Market Strategy"
    text = _truncate(str(topic).strip(), 90)
    return text if text else "Flexworks Market Intelligence Presentation"


def _speaker_note_for(slide: Mapping[str, Any]) -> str:
    bullets = slide.get("bullets") or []
    point = slide.get("takeaway") or "Use this slide to move the story forward."
    return "Talk through the slide as a concise decision-support point. " + point + (" Key bullets: " + "; ".join(bullets[:3]) if bullets else "")


def _add_slide_body(slide: Any, slide_data: PresentationSlide, colors: Mapping[str, Any], MSO_SHAPE: Any, Inches: Any, Pt: Any) -> None:
    left_width = 7.1 if slide_data.visualType != "none" else 11.3
    _add_bullets(slide, slide_data.bullets, 0.75, 1.75, left_width, 2.85, Pt(15), colors)
    if slide_data.takeaway:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(5.35), Inches(left_width), Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = colors["white"]
        box.line.color.rgb = colors["border"]
        text_frame = box.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = slide_data.takeaway
        paragraph.font.size = Pt(11)
        paragraph.font.bold = True
        paragraph.font.color.rgb = colors["dark"]
    if slide_data.visualType != "none":
        visual = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.45), Inches(1.65), Inches(4.25), Inches(3.95))
        visual.fill.solid()
        visual.fill.fore_color.rgb = colors["white"]
        visual.line.color.rgb = colors["border"]
        text_frame = visual.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = slide_data.visualType.replace("_", " ").title()
        paragraph.font.size = Pt(16)
        paragraph.font.bold = True
        paragraph.font.color.rgb = colors["accent"]
        if slide_data.visualPrompt:
            note = text_frame.add_paragraph()
            note.text = slide_data.visualPrompt
            note.font.size = Pt(10)
            note.font.color.rgb = colors["muted"]


def _add_bullets(slide: Any, bullets: Sequence[str], x: float, y: float, w: float, h: float, font_size: Any, colors: Mapping[str, Any]) -> None:
    shape = slide.shapes.add_textbox(_inches(x), _inches(y), _inches(w), _inches(h))
    text_frame = shape.text_frame
    text_frame.clear()
    for index, bullet in enumerate(bullets[:5]):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = font_size
        paragraph.font.color.rgb = colors["dark"]
        paragraph.space_after = font_size


def _add_textbox(slide: Any, text: str, x: float, y: float, w: float, h: float, font_size: Any, color: Any, *, bold: bool = False) -> None:
    textbox = slide.shapes.add_textbox(_inches(x), _inches(y), _inches(w), _inches(h))
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = font_size
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def _add_accent_bar(slide: Any, x: float, y: float, w: float, h: float, color: Any) -> None:
    shape = slide.shapes.add_shape(1, _inches(x), _inches(y), _inches(w), _inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def _add_footer(slide: Any, deck_title: str, slide_number: int, slide_count: int, colors: Mapping[str, Any], Inches: Any, Pt: Any) -> None:
    _add_textbox(slide, _truncate(deck_title, 70), 0.55, 7.02, 9.5, 0.25, Pt(8), colors["muted"])
    footer = slide.shapes.add_textbox(Inches(11.7), Inches(7.02), Inches(1.0), Inches(0.25))
    paragraph = footer.text_frame.paragraphs[0]
    paragraph.text = f"{slide_number}/{slide_count}"
    paragraph.alignment = 2
    paragraph.font.size = Pt(8)
    paragraph.font.color.rgb = colors["muted"]


def _add_speaker_notes(slide: Any, notes: str) -> None:
    try:
        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.text = notes
    except Exception:
        return


def _inches(value: float) -> int:
    from pptx.util import Inches

    return Inches(value)


def _choice(value: object, options: Mapping[str, str], default: str) -> str:
    text = str(value or "").strip()
    if text in options:
        return text
    reverse = {label: key for key, label in options.items()}
    return reverse.get(text, default)


def _source_material(normalized_input: Mapping[str, Any]) -> str:
    if normalized_input.get("sourceType") == "topic":
        return _clean_text(normalized_input.get("topic")) or _clean_text(normalized_input.get("sourceText"))
    return _clean_text(normalized_input.get("sourceText")) or _clean_text(normalized_input.get("topic"))


def _required_text(value: object, field_name: str) -> str:
    text = _clean_text(value)
    if not text:
        raise ValueError(f"Missing required presentation field: {field_name}.")
    return text


def _optional_text(value: object, max_chars: int) -> str | None:
    text = _clean_text(value)
    return _truncate(text, max_chars) if text else None


def _clean_text(value: object) -> str:
    if value is None:
        return ""
    return re.sub(r"\s+", " ", str(value)).strip()


def _plain_text(value: str) -> str:
    return re.sub(r"[#*_`>\[\]()]", "", value or "").strip()


def _truncate(text: str, max_chars: int) -> str:
    text = _clean_text(text)
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "…"


def _list_of_text(value: object, *, max_items: int, max_chars: int) -> list[str]:
    if value is None:
        return []
    if not isinstance(value, Sequence) or isinstance(value, (str, bytes)):
        return []
    return [_truncate(_clean_text(item), max_chars) for item in value if _clean_text(item)][:max_items]


def _chart_rows(value: object) -> list[dict[str, str | int | float]]:
    if not isinstance(value, Sequence) or isinstance(value, (str, bytes)):
        return []
    rows = []
    for row in value[:20]:
        if isinstance(row, Mapping):
            cleaned = {str(key): cell for key, cell in row.items() if isinstance(cell, (str, int, float))}
            if cleaned:
                rows.append(cleaned)
    return rows


def _table_rows(value: object) -> list[list[str]]:
    if not isinstance(value, Sequence) or isinstance(value, (str, bytes)):
        return []
    rows = []
    for row in value[:12]:
        if isinstance(row, Sequence) and not isinstance(row, (str, bytes)):
            rows.append([_truncate(_clean_text(cell), 120) for cell in row[:6]])
    return rows


def _infer_title_from_source(source: str) -> str:
    for line in source.splitlines():
        clean = _plain_text(line)
        if clean:
            return _truncate(clean, 90)
    return ""


def _label(metric_column: str) -> str:
    return {
        "Selected_Metric": "Selected metric",
        "Revenue_per_kW": "Revenue per kW",
        "Annualized_Revenue": "Annualized Revenue",
        "Opportunity_Score": "Opportunity Score",
        "Risk_Adjusted_Score": "Risk-adjusted Score",
    }.get(metric_column, metric_column.replace("_", " "))


def _format_value(value: float, metric_column: str) -> str:
    if "Revenue_per_kW" in metric_column:
        return f"${value:,.2f}/kW"
    if "Revenue" in metric_column:
        return f"${value:,.0f}"
    return f"{value:,.2f}"
